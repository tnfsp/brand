# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import os
from datetime import datetime

# ==================== 貼文資訊配置 ====================
POST_CONFIG = {
    'date': '2025-12-12',
    'type': '個人成長',  # 醫學 / 故事 / 個人成長
    'title': '沒有姓名的值班機器',
    'subtitle': 'PART 1',  # 副標題（成長主題專用）
}

# 模板檔案映射
TEMPLATE_MAP = {
    '醫學': 'assets/canva-templates/醫學貼文模板.pptx',
    '故事': 'assets/canva-templates/故事貼文模板 .pptx',
    '個人成長': 'assets/canva-templates/成長貼文模板 .pptx',
}

# 自動生成資料夾路徑
folder_name = f"{POST_CONFIG['date']}-IG貼文-{POST_CONFIG['type']}-{POST_CONFIG['title']}"
output_dir = os.path.join('outputs', 'IG貼文', folder_name)

# 創建資料夾（如果不存在）
os.makedirs(output_dir, exist_ok=True)

# 檔案路徑
ppt_file = os.path.join(output_dir, f"{POST_CONFIG['title']}_草稿.pptx")

# 從模板創建簡報
template_path = TEMPLATE_MAP[POST_CONFIG['type']]
prs = Presentation(template_path)
print(f'Loaded template: {template_path}')
print(f'Template has {len(prs.slides)} slides')

# 顏色定義
BEIGE_BG = RGBColor(233, 216, 166)  # #E9D8A6
BLACK_TEXT = RGBColor(0, 0, 0)
ORANGE = RGBColor(202, 103, 2)  # #CA6702
WHITE = RGBColor(255, 255, 255)

# 浮水印圖片路徑與精確位置
WATERMARK_PATH = 'assets/reference/IG某某醫師浮水印(放在右下角).png'

def add_watermark(slide):
    """在投影片右下角添加浮水印圖片"""
    if os.path.exists(WATERMARK_PATH):
        # 浮水印精確尺寸與位置
        left = Cm(19.05)   # 19.05cm 從左上角
        top = Cm(22.35)    # 22.35cm 從左上角
        width = Cm(6.49)   # 6.49cm 寬度
        height = Cm(2.03)  # 2.03cm 高度
        slide.shapes.add_picture(WATERMARK_PATH, left, top, width=width, height=height)
    else:
        print('Warning: Watermark image not found')

# ==================== 第一張（模板封面）不修改 ====================
# 重要：第一張封面完全保留模板原樣，不做任何修改
print('Keeping slide 1 (template cover) unchanged')

# ==================== 刪除模板中間的範例內容 ====================
# 保留第一張（index 0）和最後一張，刪除中間的
# 從後往前刪除避免 index 問題
slides_to_delete = list(range(1, len(prs.slides) - 1))
for idx in reversed(slides_to_delete):
    rId = prs.slides._sldIdLst[idx].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]

print(f'Deleted example slides, now has {len(prs.slides)} slides')

# ==================== 準備插入新內容頁 ====================
# 現在 prs 只有 2 張：第一張（模板封面）和最後一張（封底）
# 我們要在第一張後面插入：
# 1. 標題頁（原本我以為的封面）
# 2. 6 張內容頁

# 先插入標題頁（第 2 張）
title_slide = prs.slides.add_slide(blank_slide_layout)

# 添加橘色背景
bg_title = title_slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg_title.fill.solid()
bg_title.fill.fore_color.rgb = ORANGE

# 標題框（黑邊 + 米黃色底）
title_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(1.5))
text_frame = title_box.text_frame
text_frame.text = POST_CONFIG['title']
text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
text_frame.paragraphs[0].font.size = Pt(60)
text_frame.paragraphs[0].font.bold = True
text_frame.paragraphs[0].font.name = 'csong'
text_frame.paragraphs[0].font.color.rgb = BLACK_TEXT
title_box.fill.solid()
title_box.fill.fore_color.rgb = BEIGE_BG
title_box.line.color.rgb = BLACK_TEXT
title_box.line.width = Pt(3)

# 副標題（成長主題專用）
if POST_CONFIG.get('subtitle'):
    subtitle = title_slide.shapes.add_textbox(Inches(2.5), Inches(5.7), Inches(5), Inches(0.6))
    sub_frame = subtitle.text_frame
    sub_frame.text = POST_CONFIG['subtitle']
    sub_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    sub_frame.paragraphs[0].font.size = Pt(32)
    sub_frame.paragraphs[0].font.name = 'csong'
    sub_frame.paragraphs[0].font.color.rgb = BLACK_TEXT

# 添加浮水印
add_watermark(title_slide)

# 註記（提醒替換 AI 圖片）
note = title_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
note_frame = note.text_frame
note_frame.text = '[請替換此投影片背景為 AI 生成的圖片]\n然後加上橘色半透明遮罩（#CA6702, 70-75% 透明度）'
note_frame.paragraphs[0].font.size = Pt(20)
note_frame.paragraphs[0].font.name = 'csong'
note_frame.paragraphs[0].font.color.rgb = WHITE

print('Added slide 2 (title page)')

# 接著插入 6 張內容頁
content_slides = [
    {
        'text': '「喂，我是某某樓層的值班醫師...」',
        'alignment': PP_ALIGN.CENTER,
        'size': 55,
    },
    {
        'lines': [
            {'text': '不是張醫師', 'size': 40},
            {'text': '不是李醫師', 'size': 40},
            {'text': '不是王醫師', 'size': 40},
            {'text': '', 'size': 40},
            {'text': '只是「值班醫師」', 'size': 50, 'bold': True},
        ],
        'alignment': PP_ALIGN.CENTER,
    },
    {
        'lines': [
            {'text': '什麼時候開始', 'size': 45, 'bold': True},
            {'text': '我們連名字都不說了？', 'size': 45, 'bold': True},
            {'text': '', 'size': 32},
            {'text': '接到學弟妹的會診電話', 'size': 32},
            {'text': '越來越常聽到這樣的開場', 'size': 32},
            {'text': '職位代替了姓名', 'size': 32},
            {'text': '功能取代了身份', 'size': 32},
        ],
        'alignment': PP_ALIGN.LEFT,
    },
    {
        'lines': [
            {'text': '為什麼會這樣？', 'size': 45, 'bold': True},
            {'text': '', 'size': 32},
            {'text': '也許是...', 'size': 32},
            {'text': '', 'size': 32},
            {'text': '工作壓力太大', 'size': 32},
            {'text': '慢慢麻木了', 'size': 32},
            {'text': '', 'size': 32},
            {'text': '害怕被記住', 'size': 32},
            {'text': '避免被針對', 'size': 32},
            {'text': '', 'size': 32},
            {'text': '或是醫療文化告訴我們', 'size': 32},
            {'text': '職位比人更重要', 'size': 32},
        ],
        'alignment': PP_ALIGN.LEFT,
    },
    {
        'lines': [
            {'text': '當我們失去名字', 'size': 45, 'bold': True},
            {'text': '我們失去了什麼？', 'size': 45, 'bold': True},
            {'text': '', 'size': 32},
            {'text': '從「人」變成「功能」', 'size': 32},
            {'text': '從「張醫師」變成「值班機器」', 'size': 32},
            {'text': '', 'size': 32},
            {'text': '這是保護機制', 'size': 32},
            {'text': '還是我們付出的代價？', 'size': 32},
        ],
        'alignment': PP_ALIGN.CENTER,
    },
    {
        'lines': [
            {'text': '你不只是一台值班機器', 'size': 50, 'bold': True},
            {'text': '', 'size': 38},
            {'text': '記得說出你的名字', 'size': 38, 'bold': True},
            {'text': '', 'size': 30},
            {'text': '下次打電話時', 'size': 30},
            {'text': '介紹你自己', 'size': 30},
            {'text': '', 'size': 30},
            {'text': '在體制中', 'size': 30},
            {'text': '不失去自己', 'size': 30},
        ],
        'alignment': PP_ALIGN.CENTER,
    },
]

# 插入內容頁
blank_slide_layout = prs.slide_layouts[6]  # 空白版面

for i, content in enumerate(content_slides, start=2):
    # 添加新投影片（會添加到最後）
    new_slide = prs.slides.add_slide(blank_slide_layout)

    # 添加米黃色背景
    background = new_slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BEIGE_BG

    # 添加文字
    if 'text' in content:
        # 單行文字
        text_box = new_slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(2))
        tf = text_box.text_frame
        tf.text = content['text']
        tf.paragraphs[0].alignment = content['alignment']
        tf.paragraphs[0].font.size = Pt(content['size'])
        tf.paragraphs[0].font.name = 'csong'
        tf.paragraphs[0].font.color.rgb = BLACK_TEXT
    elif 'lines' in content:
        # 多行文字
        text_box = new_slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(7))
        tf = text_box.text_frame

        for j, line_config in enumerate(content['lines']):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = line_config['text']
            p.alignment = content['alignment']
            p.font.size = Pt(line_config['size'])
            p.font.name = 'csong'
            p.font.color.rgb = BLACK_TEXT
            if line_config.get('bold', False):
                p.font.bold = True

    # 添加浮水印
    add_watermark(new_slide)

    print(f'Added content slide {i}')

print(f'Total slides after adding content: {len(prs.slides)}')

# ==================== 重新排序投影片 ====================
# 現在順序是：模板封面、封底、標題頁、內容1-6
# 需要改成：模板封面、標題頁、內容1-6、封底

# 移動封底（index 1）到最後
slides_list = list(prs.slides._sldIdLst)
back_cover = slides_list.pop(1)  # 取出封底（index 1）
slides_list.append(back_cover)   # 添加到最後

# 重新構建投影片列表
prs.slides._sldIdLst[:] = slides_list

print('Reordered slides: Template cover -> Title page -> Content 1-6 -> Back cover')

# ==================== 最後一張（封底）已經存在，不需要修改 ====================

# 儲存
prs.save(ppt_file)
print('PPT generation completed successfully!')
print(f'File saved: {ppt_file}')
print(f'Folder: {output_dir}')
print('')
print('Next steps:')
print('1. Generate AI cover image (bear as main character)')
print('2. Import PPT to Canva')
print('3. Replace slide 1 background with AI image')
print('4. Export and publish')
